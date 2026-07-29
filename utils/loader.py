from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document

from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd
import os


def load_document(path):

    extension = os.path.splitext(path)[1].lower()

    # -----------------------------
    # PDF
    # -----------------------------
    if extension == ".pdf":

        loader = PyPDFLoader(path)

        return loader.load()

    # -----------------------------
    # Word
    # -----------------------------
    elif extension == ".docx":

        doc = DocxDocument(path)

        text = "\n".join(
            paragraph.text
            for paragraph in doc.paragraphs
        )

        return [
            Document(
                page_content=text,
                metadata={
                    "source": path
                }
            )
        ]

    # -----------------------------
    # TXT
    # -----------------------------
    elif extension == ".txt":

        with open(path, "r", encoding="utf-8") as f:

            text = f.read()

        return [
            Document(
                page_content=text,
                metadata={
                    "source": path
                }
            )
        ]

    # -----------------------------
    # CSV
    # -----------------------------
    elif extension == ".csv":

        df = pd.read_csv(path)

        text = df.to_string(index=False)

        return [
            Document(
                page_content=text,
                metadata={
                    "source": path
                }
            )
        ]

    # -----------------------------
    # PowerPoint
    # -----------------------------
    elif extension == ".pptx":

        prs = Presentation(path)

        text = ""

        for slide in prs.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text += shape.text + "\n"

        return [
            Document(
                page_content=text,
                metadata={
                    "source": path
                }
            )
        ]

    else:

        raise ValueError(
            f"Unsupported file type: {extension}"
        )